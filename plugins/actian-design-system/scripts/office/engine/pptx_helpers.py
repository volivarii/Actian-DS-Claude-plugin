"""Actian brand helpers for the hybrid editing method (see editing.md).

Add slides from the branded template, then draw real visuals on top — the slides keep
the master footer (page number, "Copyright © Actian Corporation", teal corner accent,
ACTIAN color logo). Import everything:

    from scripts.office.engine.pptx_helpers import *      # run from the skill root

Brand tokens and chart styling come from brand.md. Every snippet here is tested.
"""
import os
import zipfile
from collections import Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ---- brand tokens (brand.md) ----
BLUE   = RGBColor(0x0E, 0x5F, 0xDB)   # primary
SKY    = RGBColor(0x3B, 0x90, 0xFE)   # secondary / links
TEAL   = RGBColor(0x03, 0xC2, 0xCD)   # accent (use ONCE per slide)
POWDER = RGBColor(0xC4, 0xE0, 0xFA)
ICE    = RGBColor(0xAA, 0xFF, 0xFF)
NAVY   = RGBColor(0x00, 0x00, 0x32)
SLATE  = RGBColor(0x39, 0x42, 0x47)   # muted text / axis labels
GREY   = RGBColor(0xF1, 0xF1, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
MIDBLUE = RGBColor(0x8A, 0xC6, 0xF8)
GRIDLINE = RGBColor(0xDC, 0xE6, 0xF0)
CHART_SERIES = [BLUE, TEAL, SKY, POWDER, MIDBLUE, NAVY, SLATE, ICE]

# layout indices (0-based; file index N == N-1). See layouts.md.
LY_TITLE, LY_TITLE_ONLY, LY_BLANK = 6, 29, 26
LY_CONTENT_TITLE_ONLY, LY_DIVIDER, LY_END = 21, 40, 52


_PLUGIN_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", ".."))
WORKING_TEMPLATE = os.path.join(_PLUGIN_ROOT, "assets", "office", "Actian-Template-2026.pptx")


def open_template(path=WORKING_TEMPLATE):
    """Path to the bundled template — a 0-slide .pptx with all 54 layouts.

    Just `Presentation(open_template())`, add slides, then save elsewhere. python-pptx reads it into
    memory and writes a new file, so the bundled template stays pristine — no copy or conversion."""
    return path


def branded_canvas(prs, title, layout=LY_TITLE_ONLY):
    """Add a slide from a branded layout and set its Arial-bold title.
    The slide inherits the master footer/logo/teal-corner. Draw visuals below the title."""
    s = prs.slides.add_slide(prs.slide_layouts[layout])
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            r = ph.text_frame.paragraphs[0].runs[0]
            r.font.name = "Arial"; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    return s


def text_box(slide, l, t, w, h, runs, size=14, color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=None):
    """Add a zero-margin Arial text box. `runs` is a string or [(text, color, bold), ...]."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str): runs = [(runs, color, bold)]
    p = tf.paragraphs[0]; p.alignment = align
    if line: p.line_spacing = line
    for txt, clr, bb in runs:
        r = p.add_run(); r.text = txt; r.font.name = "Arial"; r.font.size = Pt(size)
        r.font.bold = bb; r.font.color.rgb = clr
    return tb


def flat_shape(slide, shp_type, l, t, w, h, fill, line_clr=None):
    """Add a flat (no-shadow) filled shape — honors the brand's flat aesthetic."""
    sh = slide.shapes.add_shape(shp_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_clr is None: sh.line.fill.background()
    else: sh.line.color.rgb = line_clr; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def style_chart(chart):
    """Apply Actian series colors + flat axis styling to a chart added via add_chart()."""
    chart.font.name = "Arial"; chart.font.size = Pt(10); chart.font.color.rgb = SLATE
    plot = chart.plots[0]
    for i, s in enumerate(plot.series):
        s.format.fill.solid(); s.format.fill.fore_color.rgb = CHART_SERIES[i % len(CHART_SERIES)]
        s.format.line.fill.background()
    cat, val = chart.category_axis, chart.value_axis
    cat.has_major_gridlines = False
    val.has_major_gridlines = True
    val.major_gridlines.format.line.color.rgb = GRIDLINE
    val.major_gridlines.format.line.width = Pt(0.5)
    for ax in (cat, val):
        ax.tick_labels.font.name = "Arial"; ax.tick_labels.font.size = Pt(10)
        ax.tick_labels.font.color.rgb = SLATE
        ax.format.line.color.rgb = GRIDLINE
    return chart


def save_dedup(prs, path):
    """Save, then re-save to dedupe zip entries (see the duplicate-entry gotcha in editing.md)."""
    prs.save(path)
    Presentation(path).save(path)
    dups = [n for n, c in Counter(zipfile.ZipFile(path).namelist()).items() if c > 1]
    assert not dups, f"duplicate zip entries remain: {dups}"
    return path
