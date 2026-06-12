"""slide-data.json -> branded .pptx deck (text slide types; PR1 scope).

Maps the 5 logical slide types onto corporate layouts and fills placeholders by
idx (robust to non-standard type assignments in the Actian template). Visual/chart
hybrid-drawing from content[] is a follow-up (see plan Task 5 note).

IMPORTANT — template findings (inspected 2026-06-12):
  layout[6]  1_Title:
    idx=10 BODY  → presentation title
    idx=14 BODY  → subtitle
    idx=15 BODY  → presenter ("Presenter Name, Presenter Title / Month #, Year") — must fill or remove
    idx=2  DATE  → "2026" — remove
  layout[18] 1_Content:
    idx=0  TITLE  → page title
    idx=11 BODY   → subtitle ("Subtitle Placeholder") — remove if unused
    idx=16 OBJECT → body ("Click to edit Master text styles") — fill with body text
    idx=17 DATE   → "2026" — remove
  layout[40] 1_Divider:
    idx=10 BODY     → divider title ("Divider Title Placeholder")
    idx=1  SUBTITLE → subtitle
    idx=11 DATE     → "2026" — remove
  layout[52] 1_End Slide (back-cover — branded gradient closer; zero placeholders):
    title and subtitle drawn as white text boxes above the centered Actian logo

No TITLE-type placeholder exists on the Title cover or Divider layouts;
all text is idx-addressed as BODY/SUBTITLE.
"""
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))
from engine.pptx_helpers import open_template, save_dedup, text_box, WHITE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

# Layout 0-based indices (file index N => python-pptx slide_layouts[N-1])
# See references/office/layouts.md
LAYOUT = {
    "cover":             6,   # 1_Title  (all BODY placeholders by idx)
    "section":          40,   # 1_Divider
    "body-full":        18,   # 1_Content
    "body-text-visual": 18,   # 1_Content (PR1: text only; Two Content upgrade deferred)
    "back-cover":       52,   # 1_End Slide (branded gradient closer; 0 placeholders — title drawn)
}


def _set_text(ph, text):
    """Write text into a placeholder and enforce Arial on every run."""
    ph.text = text
    for p in ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"


def _remove_ph(ph):
    """Remove a placeholder element from its parent shape tree."""
    sp = ph._element
    sp.getparent().remove(sp)


def _ph_by_idx(slide, idx):
    """Return the placeholder with the given idx, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _fill_cover(slide, slide_data):
    """
    1_Title layout — all non-TITLE placeholders, addressed by idx:
      idx=10 BODY  → title
      idx=14 BODY  → subtitle
      idx=15 BODY  → meta (date · creators); remove if empty
      idx=2  DATE  → remove
    """
    title = slide_data.get("title") or ""
    subtitle = slide_data.get("subtitle") or ""
    meta = " · ".join(x for x in (slide_data.get("date"), slide_data.get("creators")) if x)

    for idx, value in ((10, title), (14, subtitle), (15, meta)):
        ph = _ph_by_idx(slide, idx)
        if ph is None:
            continue
        if value:
            _set_text(ph, value)
        else:
            _remove_ph(ph)

    # Remove date placeholder — contains "2026" but we don't want to display it
    date_ph = _ph_by_idx(slide, 2)
    if date_ph is not None:
        _remove_ph(date_ph)


def _fill_section(slide, slide_data):
    """
    1_Divider layout:
      idx=10 BODY     → title ("Divider Title Placeholder")
      idx=1  SUBTITLE → subtitle
      idx=11 DATE     → remove
    """
    title = slide_data.get("title") or ""
    subtitle = slide_data.get("subtitle") or ""

    title_ph = _ph_by_idx(slide, 10)
    if title_ph is not None:
        if title:
            _set_text(title_ph, title)
        else:
            _remove_ph(title_ph)

    subtitle_ph = _ph_by_idx(slide, 1)
    if subtitle_ph is not None:
        if subtitle:
            _set_text(subtitle_ph, subtitle)
        else:
            _remove_ph(subtitle_ph)

    date_ph = _ph_by_idx(slide, 11)
    if date_ph is not None:
        _remove_ph(date_ph)


def _fill_body(slide, slide_data):
    """
    1_Content layout:
      idx=0  TITLE  → title
      idx=16 OBJECT → body text (PR1: plain text; chart deferred)
      idx=11 BODY   → subtitle — remove (leftover "Subtitle Placeholder")
      idx=17 DATE   → remove
    """
    title = slide_data.get("title") or ""
    body = slide_data.get("body") or ""

    title_ph = _ph_by_idx(slide, 0)
    if title_ph is not None:
        if title:
            _set_text(title_ph, title)
        else:
            _remove_ph(title_ph)

    body_ph = _ph_by_idx(slide, 16)
    if body_ph is not None:
        if body:
            _set_text(body_ph, body)
        else:
            _remove_ph(body_ph)

    # Subtitle placeholder: always remove — we don't use it and it has prompt text
    subtitle_ph = _ph_by_idx(slide, 11)
    if subtitle_ph is not None:
        _remove_ph(subtitle_ph)

    date_ph = _ph_by_idx(slide, 17)
    if date_ph is not None:
        _remove_ph(date_ph)


def _fill_back_cover(slide, slide_data):
    """
    1_End Slide layout (layout[52]) — zero fillable placeholders.

    The corporate gradient background, centered Actian logo, "a division of
    HCLSoftware" lockup, and actian.com link are all STATIC shapes in the
    layout and render automatically.  We draw the title and optional subtitle
    as white centered text boxes in the clear space above the logo.

    Layout 52 shape geometry (inspected 2026-06-12):
      Picture 7 (Actian logo):           top=2.66in, h=1.32in (bottom ~3.98in)
      Picture 9 (HCLSoftware lockup):    top=4.65in
      TextBox 4 (actian.com):            top=6.48in
      Rectangle 5 (background):          full-bleed gradient

    Clear space above the logo: ~0.34in – 2.66in.
    Title is placed at top=1.0in; subtitle (if present) at top=2.1in,
    giving ≥0.56in clearance below the subtitle before the logo.
    """
    title = slide_data.get("title") or "Thank you"
    text_box(slide, 0.64, 1.0, 12.05, 1.4, title,
             size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    subtitle = slide_data.get("subtitle")
    if subtitle:
        text_box(slide, 0.64, 2.1, 12.05, 1.0, subtitle,
                 size=18, color=WHITE, align=PP_ALIGN.CENTER)


def render_presentation(data, out_path):
    """Render a slide-data.json dict into a branded .pptx at out_path.

    Returns out_path on success.
    """
    prs = Presentation(open_template())

    for slide_data in data.get("slides", []):
        stype = slide_data.get("type")
        layout_idx = LAYOUT.get(stype, LAYOUT["body-full"])
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        if stype == "cover":
            _fill_cover(slide, slide_data)
        elif stype == "section":
            _fill_section(slide, slide_data)
        elif stype in ("body-full", "body-text-visual"):
            _fill_body(slide, slide_data)
        elif stype == "back-cover":
            _fill_back_cover(slide, slide_data)
        else:
            # Unknown type: remove all placeholders to avoid prompt leakage
            for ph in list(slide.placeholders):
                _remove_ph(ph)

    return save_dedup(prs, out_path)
