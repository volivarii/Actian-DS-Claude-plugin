import os

from pptx import Presentation

from mappers.presentation_pptx import render_presentation
from engine.qa import PROMPT_MARKERS


def _all_text(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
    return "\n".join(out).lower()


def test_renders_one_slide_per_input(tmp_path, fixture_slide_data):
    out = tmp_path / "deck.pptx"
    render_presentation(fixture_slide_data, str(out))
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == len(fixture_slide_data["slides"])  # 5


def test_titles_present_and_no_leftover_prompts(tmp_path, fixture_slide_data):
    out = tmp_path / "deck.pptx"
    render_presentation(fixture_slide_data, str(out))
    prs = Presentation(str(out))
    text = _all_text(prs)
    assert "actian design system 2026" in text   # cover title
    assert "thank you" in text                    # back-cover title
    for marker in PROMPT_MARKERS:
        assert marker not in text, f"leftover prompt text: {marker!r}"


def test_every_run_is_arial(tmp_path, fixture_slide_data):
    out = tmp_path / "deck.pptx"
    render_presentation(fixture_slide_data, str(out))
    prs = Presentation(str(out))
    for s in prs.slides:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        assert r.font.name == "Arial", f"non-Arial run: {r.text!r}"


def test_back_cover_uses_end_slide_layout(tmp_path, fixture_slide_data):
    out = tmp_path / "deck.pptx"
    render_presentation(fixture_slide_data, str(out))
    prs = Presentation(str(out))
    assert prs.slides[-1].slide_layout.name == "1_End Slide"
